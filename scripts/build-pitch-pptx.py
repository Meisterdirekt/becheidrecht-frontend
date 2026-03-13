#!/usr/bin/env python3
"""
BescheidRecht — Vertriebs-Präsentation für soziale Einrichtungen
Generiert eine professionelle PowerPoint-Datei.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
SKY = RGBColor(0x0E, 0xA5, 0xE9)
SKY_LIGHT = RGBColor(0xE0, 0xF2, 0xFE)
SKY_DARK = RGBColor(0x07, 0x59, 0x85)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT = RGBColor(0xFE, 0xF2, 0xF2)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
GREEN = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ──
def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=SLATE_900,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri", spacing=None,
             line_spacing=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if spacing:
        p.space_before = Pt(spacing)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox

def add_multiline(slide, left, top, width, height, lines, default_size=14,
                  default_color=SLATE_500, font_name="Calibri", line_spacing=None,
                  align=PP_ALIGN.LEFT):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox

def slide_number(slide, num, total=10, color=SLATE_400):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num} / {total}", size=9, color=color, align=PP_ALIGN.RIGHT,
             bold=True)

def add_accent_bar(slide, top=0):
    add_rect(slide, 0, top, SLIDE_W, Pt(4), fill=SKY)

def section_label(slide, text, left=Inches(0.8), top=Inches(0.8)):
    add_text(slide, left, top, Inches(6), Inches(0.3), text,
             size=10, color=SKY, bold=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_accent_bar(sl)

# Subtle gradient bar at top
add_rect(sl, 0, Pt(4), SLIDE_W, Inches(0.6), fill=SKY_LIGHT)

section_label(sl, "PARTNER-PRÄSENTATION  ·  VERTRAULICH  ·  2026", top=Inches(1.4))

# Title
add_multiline(sl, Inches(0.8), Inches(2.0), Inches(8), Inches(1.6), [
    ("BescheidRecht", 54, SLATE_900, True),
])

# Subtitle
add_multiline(sl, Inches(0.8), Inches(3.3), Inches(7), Inches(1.2), [
    ("Technische Bescheid-Analyse", 26, SLATE_700, True),
    ("für soziale Einrichtungen.", 26, SLATE_700, True),
], line_spacing=34)

# Description
add_text(sl, Inches(0.8), Inches(4.6), Inches(6.5), Inches(0.8),
         "Wie Einrichtungen wie Caritas, AWO und Diakonie täglich Stunden sparen —\nund mehr Klientinnen und Klienten strukturiert unterstützen.",
         size=14, color=SLATE_500, line_spacing=22)

# Partner pills
partners = ["Caritas", "AWO", "Diakonie", "DRK", "Paritätischer"]
x_start = Inches(0.8)
for i, p in enumerate(partners):
    pill = add_rounded_rect(sl, x_start + Inches(i * 1.6), Inches(5.8),
                            Inches(1.4), Inches(0.4), fill=SLATE_100, border_color=SLATE_200)
    add_text(sl, x_start + Inches(i * 1.6) + Inches(0.1), Inches(5.83),
             Inches(1.2), Inches(0.35), p, size=11, color=SLATE_700, bold=True,
             align=PP_ALIGN.CENTER)

# Stats bar at bottom
stats = [("163+", "Fehlertypen"), ("16", "Rechtsgebiete"), ("< 90s", "Analysezeit"), ("DSGVO", "Art. 25")]
add_rect(sl, 0, Inches(6.6), SLIDE_W, Inches(0.9), fill=WHITE)
add_rect(sl, 0, Inches(6.55), SLIDE_W, Pt(1), fill=SLATE_200)
for i, (val, lbl) in enumerate(stats):
    cx = Inches(2.0 + i * 2.8)
    add_text(sl, cx, Inches(6.6), Inches(1.5), Inches(0.35), val, size=22, color=SKY,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, cx, Inches(6.95), Inches(1.5), Inches(0.25), lbl, size=8, color=SLATE_400,
             bold=True, align=PP_ALIGN.CENTER)

# Disclaimer
add_text(sl, Inches(0.8), Inches(7.1), Inches(6), Inches(0.3),
         "Technisches Analyse-Werkzeug gem. § 2 Abs. 1 RDG. Kein Ersatz für Rechtsberatung.",
         size=7, color=SLATE_400)

slide_number(sl, 1)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — DAS PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "01  ·  DAS PROBLEM")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
         "Jeder fehlerhafte Bescheid\nkostet Ihre Klienten bares Geld.",
         size=36, color=SLATE_900, bold=True, line_spacing=44)

add_text(sl, Inches(0.8), Inches(2.6), Inches(7), Inches(0.6),
         "Soziale Einrichtungen kämpfen täglich: zu viele Fälle, zu wenig Zeit, keine KI-Werkzeuge.",
         size=15, color=SLATE_500, line_spacing=22)

# Problem cards
problems = [
    ("42 %", "Erfolgreiche Widersprüche*", "Der Widersprüche gegen Jobcenter-Bescheide\nsind erfolgreich — die Fehler sind da,\nsie werden nur nicht systematisch gefunden.", RED, RED_LIGHT),
    ("3–8 h", "Pro komplexem Fall", "Verbringt eine Fachkraft mit manueller\nRecherche, Paragraph-Suche und Brief-\nentwurf — Zeit, die für Klienten fehlt.", AMBER, AMBER_LIGHT),
    ("30 Tage", "Widerspruchsfrist", "Laufen ohne KI-Fristüberwachung ungenutzt\nab. Verpasste Frist = kein Widerspruch\nmehr möglich.", RED, RED_LIGHT),
]

for i, (stat, label, desc, accent, bg) in enumerate(problems):
    x = Inches(0.8 + i * 3.95)
    card = add_rounded_rect(sl, x, Inches(3.5), Inches(3.7), Inches(2.8), fill=bg, border_color=accent)
    add_text(sl, x + Inches(0.4), Inches(3.7), Inches(2.9), Inches(0.6),
             stat, size=38, color=accent, bold=True)
    add_text(sl, x + Inches(0.4), Inches(4.35), Inches(2.9), Inches(0.3),
             label, size=11, color=SLATE_900, bold=True)
    add_text(sl, x + Inches(0.4), Inches(4.8), Inches(2.9), Inches(1.2),
             desc, size=10, color=SLATE_700, line_spacing=16)

# Quote
add_rounded_rect(sl, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.65), fill=WHITE, border_color=SLATE_200)
add_text(sl, Inches(1.1), Inches(6.58), Inches(9), Inches(0.3),
         "\u201eUnsere Berater:innen wissen oft, dass etwas nicht stimmt \u2014 aber sie haben nicht die Zeit, es systematisch zu pr\u00fcfen.\u201c",
         size=11, color=SLATE_700, bold=True)
add_text(sl, Inches(1.1), Inches(6.88), Inches(9), Inches(0.2),
         "Typisches Feedback aus Sozialberatungen", size=9, color=SLATE_400)

# Source
add_text(sl, Inches(0.8), Inches(7.15), Inches(11), Inches(0.2),
         "* Quelle: Statistik der Bundesagentur für Arbeit, Widersprüche und Klagen SGB II (2024)",
         size=7, color=SLATE_400)

slide_number(sl, 2)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — DIE LÖSUNG
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "02  ·  DIE LÖSUNG")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
         "BescheidRecht: Technische Analyse\nfür Ihre Fachkräfte.",
         size=36, color=SLATE_900, bold=True, line_spacing=44)

add_text(sl, Inches(0.8), Inches(2.6), Inches(7), Inches(0.6),
         "Fachkraft lädt den Bescheid hoch — BescheidRecht analysiert auf Unstimmigkeiten\nund erstellt eine Musterschreiben-Vorlage. In unter 90 Sekunden.",
         size=15, color=SLATE_500, line_spacing=22)

features = [
    ("⚡", "13 KI-Agenten in Parallel", "Spezialisierte Agenten prüfen gleichzeitig\nauf Fristen, Rechtsbasis, Formfehler\nund Begründungsmängel."),
    ("📄", "163+ geprüfte Fehlertypen", "Technische Prüfung über 16 Rechtsgebiete:\nSGB II bis XII, BAMF, BAföG,\nWohngeld und mehr."),
    ("🔒", "DSGVO by Design", "Vollautomatische Pseudonymisierung aller\nKlientendaten vor jeder KI-Analyse.\nArt. 25 DSGVO erfüllt."),
    ("⏱️", "< 90 Sekunden", "Von Hochladen bis Analyse-Ergebnis\nmit Musterschreiben-Vorlage\nals DIN A4 PDF."),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(3.5 + row * 1.9)
    card = add_rounded_rect(sl, x, y, Inches(5.9), Inches(1.7), fill=SLATE_100, border_color=SLATE_200)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.35), Inches(0.7), Inches(0.7)) if False else None
    add_text(sl, x + Inches(0.3), y + Inches(0.25), Inches(0.7), Inches(0.7),
             icon, size=24, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(1.2), y + Inches(0.25), Inches(4.2), Inches(0.35),
             title, size=14, color=SLATE_900, bold=True)
    add_text(sl, x + Inches(1.2), y + Inches(0.65), Inches(4.2), Inches(1.0),
             desc, size=11, color=SLATE_500, line_spacing=16)

slide_number(sl, 3)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — SO FUNKTIONIERT ES
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
# Full sky background
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=SKY)

add_text(sl, Inches(0.8), Inches(0.7), Inches(6), Inches(0.3),
         "03  ·  SO FUNKTIONIERT ES", size=10, color=RGBColor(0xBA, 0xE6, 0xFD), bold=True)

add_text(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
         "In 3 Schritten\nzur strukturierten Analyse.",
         size=36, color=WHITE, bold=True, line_spacing=44)

steps = [
    ("01", "Bescheid hochladen", "Fachkraft lädt PDF oder Scan hoch.\nAutomatische Pseudonymisierung\nschützt alle Klientendaten sofort —\nkein manueller Aufwand."),
    ("02", "KI analysiert", "13 spezialisierte Agenten prüfen auf\n163+ Fehlertypen, berechnen Fristen\nund identifizieren Unstimmigkeiten —\nin unter 90 Sekunden."),
    ("03", "Vorlage exportieren", "Musterschreiben-Vorlage als DIN A4\nPDF. Fachkraft prüft, ergänzt —\nals Basis für das Gespräch mit\nAnwalt oder Sozialverband."),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(sl, x, Inches(2.8), Inches(3.8), Inches(3.2),
                            fill=RGBColor(0x38, 0xBD, 0xF8), border_color=RGBColor(0x7D, 0xD3, 0xFC))
    # Step number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(3.1), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    add_text(sl, x + Inches(0.3), Inches(3.15), Inches(0.6), Inches(0.55),
             num, size=14, color=SKY, bold=True, align=PP_ALIGN.CENTER)

    add_text(sl, x + Inches(0.3), Inches(3.85), Inches(3.2), Inches(0.4),
             title, size=16, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.3), Inches(4.35), Inches(3.2), Inches(1.5),
             desc, size=11, color=RGBColor(0xE0, 0xF2, 0xFE), line_spacing=17)

# Arrow between steps
for i in range(2):
    x = Inches(4.45 + i * 4.1)
    add_text(sl, x, Inches(4.0), Inches(0.5), Inches(0.5),
             "→", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Bottom stats
stats_bottom = [("163+", "Fehlertypen"), ("16", "Rechtsgebiete"), ("< 90s", "Analysezeit"), ("100%", "DSGVO-konform")]
for i, (val, lbl) in enumerate(stats_bottom):
    cx = Inches(1.5 + i * 2.8)
    card = add_rounded_rect(sl, cx, Inches(6.3), Inches(2.2), Inches(0.9),
                            fill=RGBColor(0x38, 0xBD, 0xF8), border_color=RGBColor(0x7D, 0xD3, 0xFC))
    add_text(sl, cx, Inches(6.35), Inches(2.2), Inches(0.4), val, size=20, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, cx, Inches(6.75), Inches(2.2), Inches(0.25), lbl, size=8,
             color=RGBColor(0xBA, 0xE6, 0xFD), bold=True, align=PP_ALIGN.CENTER)

slide_number(sl, 4, color=RGBColor(0x7D, 0xD3, 0xFC))

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — ZIELGRUPPE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "04  ·  ZIELGRUPPE")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
         "Gebaut für soziale Einrichtungen.", size=36, color=SLATE_900, bold=True)

add_text(sl, Inches(0.8), Inches(2.1), Inches(7), Inches(0.6),
         "Überall dort, wo Fachkräfte täglich Bescheide prüfen — für Klienten, die strukturierte Unterstützung brauchen.",
         size=15, color=SLATE_500, line_spacing=22)

orgs = [
    ("Caritas", "Sozialberatung, Schuldnerberatung,\nMigrationsberatung", "Jobcenter · Sozialamt · BAMF"),
    ("AWO", "Allgemeine Sozialberatung,\nJugend- und Familienhilfe", "Jugendamt · Jobcenter · Rente"),
    ("Diakonie", "Wohnungslosenberatung,\nSuchtberatung, Flüchtlingshilfe", "Sozialamt · Jobcenter · BAMF"),
    ("DRK + Paritätischer", "Krankenhaussozialdienst,\nPflegeberatung, Seniorenhilfe", "Krankenkasse · Pflegekasse · MDK"),
]

for i, (org, role, cases) in enumerate(orgs):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(3.0 + row * 1.85)
    card = add_rounded_rect(sl, x, y, Inches(5.9), Inches(1.65), fill=SLATE_100, border_color=SLATE_200)
    add_text(sl, x + Inches(0.4), y + Inches(0.2), Inches(5), Inches(0.35),
             org, size=18, color=SLATE_900, bold=True)
    add_text(sl, x + Inches(0.4), y + Inches(0.6), Inches(5), Inches(0.6),
             role, size=11, color=SLATE_500, line_spacing=16)
    add_text(sl, x + Inches(0.4), y + Inches(1.2), Inches(5), Inches(0.3),
             cases, size=10, color=SKY, bold=True)

# Big number callout
add_rounded_rect(sl, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.65), fill=SKY_LIGHT, border_color=SKY)
add_text(sl, Inches(1.1), Inches(6.62), Inches(11), Inches(0.3),
         "In Deutschland gibt es über 1.400 Caritas- und AWO-Kreisverbände allein.",
         size=14, color=SLATE_900, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(1.1), Inches(6.95), Inches(11), Inches(0.2),
         "Dazu kommen Diakonie, DRK, Paritätischer und hunderte kommunale Beratungsstellen.",
         size=10, color=SLATE_500, align=PP_ALIGN.CENTER)

slide_number(sl, 5)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — ROI / WIRTSCHAFTLICHKEIT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "05  ·  WIRTSCHAFTLICHKEIT")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
         "Der ROI rechnet sich ab Tag 1.", size=36, color=SLATE_900, bold=True)

add_text(sl, Inches(0.8), Inches(2.1), Inches(7), Inches(0.6),
         "Rechenbeispiel: Einrichtung mit 5 Beratern, 20 Bescheide pro Berater/Monat, Ø 35 €/h Personalkosten.",
         size=14, color=SLATE_500, line_spacing=21)

# ROI calculation visual
roi_items = [
    ("100", "Bescheide / Monat", "5 Berater × 20 Bescheide", SLATE_100),
    ("40 h", "Zeitersparnis / Monat", "100 Bescheide × 24 Min. Ersparnis", SKY_LIGHT),
    ("1.400 €", "Wert der Zeitersparnis", "40 Stunden × 35 €/h", SKY_LIGHT),
    ("299 €", "Starter-Tarif", "bis 100 Analysen / Monat", WHITE),
]

for i, (val, label, sub, bg) in enumerate(roi_items):
    y = Inches(3.0 + i * 0.95)
    add_rounded_rect(sl, Inches(0.8), y, Inches(5.5), Inches(0.8), fill=bg, border_color=SLATE_200)
    add_text(sl, Inches(1.1), y + Inches(0.1), Inches(1.5), Inches(0.35),
             val, size=20, color=SKY if "€" in val or "h" in val else SLATE_900, bold=True)
    add_text(sl, Inches(2.8), y + Inches(0.1), Inches(3), Inches(0.25),
             label, size=12, color=SLATE_900, bold=True)
    add_text(sl, Inches(2.8), y + Inches(0.4), Inches(3), Inches(0.25),
             sub, size=10, color=SLATE_400)

# Big ROI result on the right
add_rounded_rect(sl, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.8), fill=RGBColor(0xF0, 0xFD, 0xF4), border_color=GREEN)
add_text(sl, Inches(7.0), Inches(3.3), Inches(5.5), Inches(0.3),
         "NETTOVORTEIL", size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(7.0), Inches(3.8), Inches(5.5), Inches(0.8),
         "+1.101 €", size=48, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(7.0), Inches(4.6), Inches(5.5), Inches(0.3),
         "pro Monat", size=16, color=SLATE_500, align=PP_ALIGN.CENTER)

add_text(sl, Inches(7.0), Inches(5.3), Inches(5.5), Inches(0.3),
         "ROI", size=10, color=SKY, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.5),
         "+368 %", size=36, color=SKY, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, Inches(7.3), Inches(6.3), Inches(5), Inches(0.3),
         "Berechnung: 25 Min. manuelle Prüfung – 1 Min. mit KI = 24 Min. Ersparnis/Bescheid",
         size=8, color=SLATE_400, align=PP_ALIGN.CENTER)

slide_number(sl, 6)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — DATENSCHUTZ & SICHERHEIT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "06  ·  DATENSCHUTZ & SICHERHEIT")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
         "DSGVO-konform by Design.\nKeine Kompromisse bei Klientendaten.",
         size=36, color=SLATE_900, bold=True, line_spacing=44)

security_items = [
    ("🔒", "Automatische Pseudonymisierung", "Alle personenbezogenen Daten (Namen, IBAN, Geburtsdaten, Adressen)\nwerden vor der KI-Analyse vollautomatisch anonymisiert.\nKeine manuellen Schritte nötig."),
    ("🇪🇺", "EU-Hosting / Art. 25 DSGVO", "Alle Daten werden in der EU verarbeitet und gespeichert.\nDSGVO Art. 25 (Privacy by Design) ist von Anfang an implementiert.\nKein Datentransfer in Drittländer."),
    ("📋", "AVV-ready", "Auftragsverarbeitungsvertrag sofort verfügbar.\nTOM-Dokumentation nach Art. 32 DSGVO liegt vor.\nKomplett vorbereitet für Ihre Datenschutzbeauftragte."),
    ("🗑️", "Keine Datenspeicherung", "Bescheid-Inhalte werden nach Analyse automatisch gelöscht.\nKeine Langzeitspeicherung von Klientendaten.\nRevisionssicher protokolliert."),
]

for i, (icon, title, desc) in enumerate(security_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.8 + row * 2.1)
    card = add_rounded_rect(sl, x, y, Inches(5.9), Inches(1.9), fill=SLATE_100, border_color=SLATE_200)
    add_text(sl, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
             icon, size=22, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.3),
             title, size=14, color=SLATE_900, bold=True)
    add_text(sl, x + Inches(1.0), y + Inches(0.55), Inches(4.5), Inches(1.2),
             desc, size=10, color=SLATE_500, line_spacing=16)

slide_number(sl, 7)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — PREISE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "07  ·  PREISE & LIZENZEN")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
         "Team-Lizenzen. Transparente Preise.", size=36, color=SLATE_900, bold=True)

add_text(sl, Inches(0.8), Inches(2.0), Inches(7), Inches(0.4),
         "Monatliche Flatrate — keine versteckten Kosten, monatlich kündbar.",
         size=15, color=SLATE_500)

tarife = [
    ("Starter", "249 €", "1 Nutzer", ["100 Analysen / Monat", "Schreiben-Generator", "Auto-Pseudonymisierung"], False),
    ("Team", "499 €", "Bis 3 Nutzer", ["400 Analysen / Monat", "Fristen-Dashboard", "Prioritäts-Support"], True),
    ("Einrichtung", "899 €", "Bis 10 Nutzer", ["1.000 Analysen / Monat", "Persönlicher Betreuer", "Prioritäts-Support"], False),
    ("Rahmenvertrag", "Auf Anfrage", "Mehrere Standorte", ["Ab 3.000 Analysen & SLA", "Compliance-Paket", "Dedizierter Betreuer"], False),
]

for i, (name, price, users, features, highlight) in enumerate(tarife):
    x = Inches(0.8 + i * 3.15)
    w = Inches(2.95)
    bg = SKY if highlight else WHITE
    border = SKY if highlight else SLATE_200
    text_primary = WHITE if highlight else SLATE_900
    text_secondary = RGBColor(0xBA, 0xE6, 0xFD) if highlight else SLATE_400
    text_feature = RGBColor(0xE0, 0xF2, 0xFE) if highlight else SLATE_500

    card = add_rounded_rect(sl, x, Inches(2.8), w, Inches(4.0), fill=bg, border_color=border)

    if highlight:
        badge = add_rounded_rect(sl, x + Inches(0.6), Inches(2.6), Inches(1.7), Inches(0.35),
                                 fill=RGBColor(0x07, 0x59, 0x85))
        add_text(sl, x + Inches(0.6), Inches(2.62), Inches(1.7), Inches(0.3),
                 "★  EMPFOHLEN", size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(sl, x + Inches(0.3), Inches(3.1), w, Inches(0.2),
             name, size=9, color=text_secondary, bold=True)
    add_text(sl, x + Inches(0.3), Inches(3.35), w, Inches(0.5),
             price, size=28, color=text_primary, bold=True)
    if "€" in price:
        add_text(sl, x + Inches(0.3), Inches(3.85), w, Inches(0.2),
                 "/ Monat", size=9, color=text_secondary)
    add_text(sl, x + Inches(0.3), Inches(4.1), w, Inches(0.2),
             users, size=10, color=SKY_LIGHT if highlight else SKY, bold=True)

    for j, feat in enumerate(features):
        add_text(sl, x + Inches(0.3), Inches(4.5 + j * 0.35), Inches(2.5), Inches(0.3),
                 f"✓  {feat}", size=9, color=text_feature)

# MwSt note
add_text(sl, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3),
         "Alle Preise zzgl. MwSt.  ·  Jahresvertrag auf Anfrage  ·  Pilotprojekte mit Sonderkonditionen möglich",
         size=9, color=SLATE_400, align=PP_ALIGN.CENTER)

slide_number(sl, 8)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — WARUM JETZT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)
section_label(sl, "08  ·  WARUM JETZT")

add_text(sl, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
         "Das Zeitfenster ist jetzt.", size=36, color=SLATE_900, bold=True)

add_text(sl, Inches(0.8), Inches(2.1), Inches(7), Inches(0.4),
         "Vier Entwicklungen treffen gerade zusammen — und schaffen ein einmaliges Momentum.",
         size=15, color=SLATE_500)

reasons = [
    ("📈", "SGB-Reformen 2025/2026", "Neue Regelsätze und Leistungsänderungen\nbedeutenungen: mehr fehlerhafte Bescheide,\nmehr Widerspruchsbedarf."),
    ("👥", "Steigende Fallzahlen", "Einrichtungen berichten von 20–40 % mehr\nBeratungsanfragen seit 2024 — ohne\nproportionale Personalaufstockung."),
    ("⚡", "KI ist produktionsreif", "2025 ist das erste Jahr, in dem KI-gestützte\nRechtsanalyse vollständig DSGVO-konform\nund praxistauglich umsetzbar ist."),
    ("🏆", "First Mover Vorteil", "Erste Einrichtungen, die einführen, setzen\nden Standard — und gewinnen Kapazitäten\nzurück, die Mitbewerber noch nicht haben."),
]

for i, (icon, title, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(3.0 + row * 1.9)
    card = add_rounded_rect(sl, x, y, Inches(5.9), Inches(1.7), fill=SLATE_100, border_color=SLATE_200)
    add_text(sl, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
             icon, size=22, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.3),
             title, size=14, color=SLATE_900, bold=True)
    add_text(sl, x + Inches(1.0), y + Inches(0.55), Inches(4.5), Inches(1.1),
             desc, size=11, color=SLATE_500, line_spacing=16)

# Bottom CTA bar
add_rounded_rect(sl, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.55), fill=SKY)
add_text(sl, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.25),
         "Einrichtungen, die jetzt starten, sparen ab Tag 1 Fachkraft-Kapazitäten.",
         size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.2),
         "Pilotprojekte mit bevorzugten Konditionen für frühe Partner",
         size=10, color=RGBColor(0xBA, 0xE6, 0xFD), align=PP_ALIGN.CENTER)

slide_number(sl, 9, color=SLATE_400)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — CTA / NÄCHSTE SCHRITTE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.8), Inches(6), Inches(0.3),
         "09  ·  NÄCHSTE SCHRITTE", size=10, color=SKY, bold=True)

add_text(sl, Inches(0), Inches(1.8), SLIDE_W, Inches(1.0),
         "Bereit für eine Demo?", size=48, color=SLATE_900, bold=True,
         align=PP_ALIGN.CENTER)

add_text(sl, Inches(3), Inches(3.0), Inches(7.3), Inches(0.8),
         "Wir zeigen Ihnen in 30 Minuten, wie BescheidRecht\nkonkret in Ihre Arbeitsabläufe passt.\nKostenlos, unverbindlich, auf Ihre Einrichtung zugeschnitten.",
         size=15, color=SLATE_500, align=PP_ALIGN.CENTER, line_spacing=24)

# CTA buttons
cta1 = add_rounded_rect(sl, Inches(3.5), Inches(4.2), Inches(3), Inches(0.65), fill=SKY)
add_text(sl, Inches(3.5), Inches(4.27), Inches(3), Inches(0.5),
         "Demo vereinbaren  →", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

cta2 = add_rounded_rect(sl, Inches(6.8), Inches(4.2), Inches(3), Inches(0.65),
                         fill=WHITE, border_color=SLATE_200)
add_text(sl, Inches(6.8), Inches(4.27), Inches(3), Inches(0.5),
         "Pilotprojekt anfragen", size=13, color=SLATE_700, bold=True, align=PP_ALIGN.CENTER)

# Quick facts
facts = [("< 1 Woche", "Onboarding"), ("Pilot", "Auf Anfrage"), ("Monatlich", "Kündbar")]
for i, (val, lbl) in enumerate(facts):
    cx = Inches(3.5 + i * 2.5)
    add_text(sl, cx, Inches(5.3), Inches(2), Inches(0.35),
             val, size=20, color=SKY, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, cx, Inches(5.65), Inches(2), Inches(0.2),
             lbl, size=9, color=SLATE_400, bold=True, align=PP_ALIGN.CENTER)

# Contact & branding
add_text(sl, Inches(0), Inches(6.3), SLIDE_W, Inches(0.4),
         "BescheidRecht", size=22, color=SLATE_900, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(6.65), SLIDE_W, Inches(0.3),
         "bescheidrecht.de  ·  kontakt@bescheidrecht.de  ·  2026",
         size=11, color=SLATE_400, align=PP_ALIGN.CENTER)

add_text(sl, Inches(0.8), Inches(7.1), Inches(6), Inches(0.3),
         "Technisches Analyse-Werkzeug gem. § 2 Abs. 1 RDG. Kein Ersatz für Rechtsberatung.",
         size=7, color=SLATE_400)

slide_number(sl, 10)

# ── Save ──
output_path = "/home/henne1990/bescheidrecht-frontend/BescheidRecht_Pitch_2026.pptx"
prs.save(output_path)
print(f"✅ Präsentation gespeichert: {output_path}")
print(f"   {len(prs.slides)} Slides · 16:9 Format · Print-ready")
